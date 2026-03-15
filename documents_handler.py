import os
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
MIN_TEXT_LENGTH = 500
import fitz  # PyMuPDF

def extract_text_from_file(file_path):
    """
    الدالة الرئيسية التي تستخرج كل شيء وتجهزه للـ Session
    """
    ext = os.path.splitext(file_path)[1].lower()
    full_text = ""
    chunks = []

    if ext == ".pdf":
        # في الـ PDF نستخرج النصوص والصور معاً
        chunks = extract_visual_chunks_from_pdf(file_path)
        full_text = "\n".join([c['text'] for c in chunks])
    
    elif ext in [".docx", ".pptx", ".txt"]:
        if ext == ".docx": full_text = extract_text_from_docx(file_path)
        elif ext == ".pptx": full_text = extract_text_from_pptx(file_path)
        else: full_text = extract_text_from_txt(file_path)
        
        full_text = clean_text(full_text)
        # تقسيم النصوص التي لا تحتوي على صور تلقائياً بناءً على الفقرات
        chunks = [{"text": chunk, "image_path": None} for chunk in full_text.split('\n\n') if len(chunk) > 20]

    if len(full_text) < MIN_TEXT_LENGTH:
        raise ValueError("المستند قصير جداً أو يحتاج لـ OCR")

    return full_text, chunks

"""
def extract_text_from_file(file_path):
    """
"""
    يحدد نوع الملف ويستخرج النص المناسب
    ويتأكد أن النص طويل كفاية للاستخدام
    """
"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        text = extract_text_from_txt(file_path)

    elif ext == ".pdf":
        return extract_visual_chunks_from_pdf(file_path)

    elif ext == ".docx":
        text = extract_text_from_docx(file_path)

    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # تنظيف النص
    text = clean_text(text)

    # التحقق من أن النص كافٍ
    if not text or len(text) < MIN_TEXT_LENGTH:
        raise ValueError("Extracted text is too short. The file may require OCR.")

    return [{"text": chunk, "image_path": None} for chunk in text.split('\n\n') if len(chunk) > 20]
    """

def extract_text_from_txt(file_path):
    """
    استخراج النص من ملفات TXT
    """
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_docx(file_path):
    """
    استخراج النص من ملفات DOCX
    """
    doc = Document(file_path)

    full_text = []

    for para in doc.paragraphs:
        full_text.append(para.text)

    return "\n".join(full_text)


def extract_text_from_pptx(file_path):
    """
    استخراج النص من ملفات PowerPoint
    """
    prs = Presentation(file_path)

    full_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)

    return "\n".join(full_text)

def extract_visual_chunks_from_pdf(file_path):
    doc = fitz.open(file_path)
    chunks = []
    
    # إنشاء مجلد مؤقت لحفظ صور الصفحات
    os.makedirs("session_images", exist_ok=True)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text").strip()
        
        # تجاهل الصفحات الفارغة جداً
        if len(text) < 20: 
            continue
            
        # تحويل الصفحة إلى صورة عالية الجودة
        img_path = f"session_images/page_{page_num}.png"
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # Matrix(2,2) لزيادة الدقة
        pix.save(img_path)
        
        chunks.append({
            "text": clean_text(text),
            "image_path": img_path
        })
        
    doc.close()
    return chunks

"""
def extract_text_from_pdf(file_path):
    
    استخراج النص من ملفات PDF
    
    reader = PdfReader(file_path)

    full_text = []

    for page in reader.pages:
        text = page.extract_text()
        if text:
            full_text.append(text)

    return "\n".join(full_text)

"""
def clean_text(text):
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)



import re

def smart_split(text, max_chars=100):
    """
    تقسيم النص إلى قطع صغيرة لضمان سرعة التوليد.
    max_chars: الحد الأقصى للحروف في القطعة الواحدة قبل كسرها قسرياً.
    """
    # 1. التنظيف الأولي
    text = text.replace('\n', ' ')
    
    # 2. التقسيم بناءً على علامات الوقف الرئيسية (. ! ؟ ?)
    # استخدمنا regex للحفاظ على علامة الوقف مع الجملة
    sentences = re.split(r'(?<=[.!?؟])\s+', text)
    
    final_chunks = []
    
    for sentence in sentences:
        if len(sentence) <= max_chars:
            final_chunks.append(sentence.strip())
        else:
            # 3. إذا كانت الجملة طويلة، نقسمها بناءً على الفواصل (, ، ;)
            sub_parts = re.split(r'(?<=[,،;؛])\s+', sentence)
            for part in sub_parts:
                if len(part) <= max_chars:
                    final_chunks.append(part.strip())
                else:
                    # 4. التقسيم القسري بناءً على عدد الكلمات (آخر حل)
                    words = part.split(' ')
                    current_chunk = []
                    current_length = 0
                    for word in words:
                        if current_length + len(word) + 1 <= max_chars:
                            current_chunk.append(word)
                            current_length += len(word) + 1
                        else:
                            final_chunks.append(" ".join(current_chunk))
                            current_chunk = [word]
                            current_length = len(word) + 1
                    if current_chunk:
                        final_chunks.append(" ".join(current_chunk))
                        
    # تصفية أي قطع فارغة
    return [c.strip() for c in final_chunks if c.strip()]

# مثال على الاستخدام:
# chunks = smart_split(long_ai_response)
# for chunk in chunks:
#     await broadcast_ai_response(chunk)
